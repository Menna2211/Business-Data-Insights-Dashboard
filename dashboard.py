import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from io import StringIO, BytesIO
import traceback
from datetime import datetime
import numpy as np
import base64
from pptx import Presentation
from pptx.util import Inches
from fpdf import FPDF
import tempfile
import os

# Configure page
st.set_page_config(
    page_title="Business Data Insights Dashboard",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Add custom CSS for better styling
st.markdown("""
<style>
    .main .block-container {
        padding-top: 2rem;
        padding-bottom: 2rem;
    }
    .stMetric {
        border-radius: 10px;
        padding: 15px;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
    }
    .stButton>button {
        width: 100%;
        border-radius: 8px;
        padding: 10px 24px;
    }
    .stRadio>div {
        display: flex;
        gap: 10px;
    }
    .stRadio>div[role="radiogroup"]>label {
        padding: 8px 12px;
        border-radius: 8px;
        border: 1px solid #e0e0e0;
    }
    .stRadio>div[role="radiogroup"]>label:hover {
        background-color: #f5f5f5;
    }
    .stRadio>div[role="radiogroup"]>label[data-baseweb="radio"]>div:first-child {
        padding-right: 8px;
    }
    .st-expander {
        border-radius: 8px;
        border: 1px solid #e0e0e0;
    }
    .stAlert {
        border-radius: 8px;
    }
    .tooltip {
        position: relative;
        display: inline-block;
    }
    .tooltip .tooltiptext {
        visibility: hidden;
        width: 200px;
        background-color: #555;
        color: #fff;
        text-align: center;
        border-radius: 6px;
        padding: 5px;
        position: absolute;
        z-index: 1;
        bottom: 125%;
        left: 50%;
        margin-left: -100px;
        opacity: 0;
        transition: opacity 0.3s;
    }
    .tooltip:hover .tooltiptext {
        visibility: visible;
        opacity: 1;
    }
    .welcome-container {
        text-align: center; 
        padding: 5rem 0;
    }
    .feature-card {
        background: #f0f2f6; 
        padding: 2rem; 
        border-radius: 10px; 
        text-align: left; 
        margin: 0 auto; 
        max-width: 800px;
    }
    .export-btn {
        margin-top: 1rem;
        margin-bottom: 1rem;
    }
</style>
""", unsafe_allow_html=True)

def detect_column_types(df):
    """Automatically detect column types (numeric, categorical, date)"""
    numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
    categorical_cols = df.select_dtypes(include=['object', 'category']).columns.tolist()
    date_cols = [col for col in df.columns if pd.api.types.is_datetime64_any_dtype(df[col])]
    
    # Try to detect date columns that weren't automatically parsed
    date_patterns = ['date', 'time', 'year', 'month', 'day']
    for col in categorical_cols:
        if any(pattern in col.lower() for pattern in date_patterns):
            try:
                df[col] = pd.to_datetime(df[col])
                date_cols.append(col)
                categorical_cols.remove(col)
            except:
                pass
    
    return {
        'numeric': numeric_cols,
        'categorical': categorical_cols,
        'date': date_cols
    }

def load_data(uploaded_file):
    """Load data with automatic format detection"""
    try:
        # Get file extension
        file_extension = uploaded_file.name.split('.')[-1].lower()
        
        # Reset file pointer
        uploaded_file.seek(0)
        
        if file_extension == 'xlsx':
            try:
                # First try with openpyxl (modern Excel files)
                df = pd.read_excel(uploaded_file, engine='openpyxl')
            except:
                try:
                    # Fall back to xlrd for older Excel files
                    df = pd.read_excel(uploaded_file, engine='xlrd')
                except Exception as e:
                    st.error(f"Failed to read Excel file: {str(e)}")
                    st.info("Please ensure: 1) File is not password protected 2) It's a valid .xlsx file 3) File is not corrupted")
                    return None
        elif file_extension == 'json':
            df = pd.read_json(uploaded_file)
        else:  # For CSV and other text files
            # Try common encodings
            encodings = ['utf-8', 'utf-8-sig', 'latin1']
            for encoding in encodings:
                try:
                    uploaded_file.seek(0)
                    content = uploaded_file.read().decode(encoding)
                    # Try common delimiters
                    for delimiter in [',', ';', '\t']:
                        try:
                            df = pd.read_csv(StringIO(content), delimiter=delimiter)
                            break
                        except:
                            continue
                    break
                except:
                    continue
        
        # Clean column names if we successfully loaded data
        if 'df' in locals():
            df.columns = [col.strip().replace(' ', '_').lower() for col in df.columns]
            return df
        else:
            st.error("Failed to read file. Please check the format and try again.")
            return None

    except Exception as e:
        st.error(f"Error loading file: {str(e)}")
        st.text(traceback.format_exc())
        return None

def create_business_summary(df):
    """Create a high-level business summary"""
    st.header("📌 Executive Summary")
    
    # Calculate basic metrics
    num_records = df.shape[0]
    num_features = df.shape[1]
    date_cols = [col for col in df.columns if pd.api.types.is_datetime64_any_dtype(df[col])]
    has_dates = len(date_cols) > 0
    
    # Create summary cards with better formatting
    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("Total Records", f"{num_records:,}", help="Number of rows in your dataset")
    with col2:
        st.metric("Total Features", num_features, help="Number of columns in your dataset")
    with col3:
        st.metric("Time Period Covered", "Available" if has_dates else "Not Available", 
                 help="Whether your data contains date/time information")
    
    # Add data quality indicators
    st.subheader("Data Quality Check")
    missing_values = df.isnull().sum().sum()
    duplicate_rows = df.duplicated().sum()
    
    col1, col2 = st.columns(2)
    with col1:
        completeness = 1 - (missing_values / (num_records * num_features))
        st.metric("Data Completeness", f"{completeness*100:.1f}%", 
                 help="Percentage of non-missing values in your dataset")
        st.progress(min(1.0, max(0.0, completeness)))
    with col2:
        uniqueness = 1 - (duplicate_rows / num_records)
        st.metric("Unique Records", f"{uniqueness*100:.1f}%", 
                 help="Percentage of unique rows in your dataset")
        st.progress(min(1.0, max(0.0, uniqueness)))
    
    # Add quick insights
    st.subheader("Quick Insights")
    if has_dates:
        date_col = date_cols[0]
        date_range = f"{df[date_col].min().strftime('%Y-%m-%d')} to {df[date_col].max().strftime('%Y-%m-%d')}"
        st.info(f"📅 **Date Range**: Your data covers **{date_range}**")
    
    numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
    if numeric_cols:
        top_num = numeric_cols[0]
        if len(df) > 1 and has_dates:
            try:
                growth = (df[top_num].iloc[-1] - df[top_num].iloc[0]) / df[top_num].iloc[0] * 100
            except:
                growth = 0
        
        st.info(f"📊 **{top_num.replace('_', ' ').title()}**:\n"
                f"- Average: {df[top_num].mean():,.2f}\n"
                f"- Maximum: {df[top_num].max():,.2f}\n"
                f"- Minimum: {df[top_num].min():,.2f}")
        
        if has_dates and len(df) > 1 and growth != 0:
            st.info(f"📈 **Growth Trend**: {'↑' if growth > 0 else '↓'} {abs(growth):.1f}% over the period")

def create_visualizations(df, col_types):
    """Generate automatic visualizations with business focus"""
    # Store all figures and data for export
    export_data = {
        'figures': [],
        'tables': {},
        'summary': {}
    }
    
    # Create summary metrics for export
    export_data['summary']['num_records'] = df.shape[0]
    export_data['summary']['num_features'] = df.shape[1]
    export_data['summary']['date_cols'] = [col for col in df.columns if pd.api.types.is_datetime64_any_dtype(df[col])]
    export_data['summary']['missing_values'] = df.isnull().sum().sum()
    export_data['summary']['duplicate_rows'] = df.duplicated().sum()
    
    create_business_summary(df)
    
    # 1. Data Overview
    with st.expander("🔍 Detailed Data Overview", expanded=False):
        st.write(f"**Shape:** {df.shape[0]} rows × {df.shape[1]} columns")
        
        tab1, tab2, tab3 = st.tabs(["📋 Sample Data", "📊 Descriptive Stats", "⚠️ Missing Values"])
        with tab1:
            st.write("First 10 rows of your data:")
            st.dataframe(df.head(10), use_container_width=True)
            export_data['tables']['sample_data'] = df.head(10)
        with tab2:
            st.write("Statistical summary of your data:")
            desc_stats = df.describe(include='all').T
            st.dataframe(desc_stats.style.background_gradient(cmap='Blues'), 
                        use_container_width=True)
            export_data['tables']['descriptive_stats'] = desc_stats
        with tab3:
            st.write("Missing values in your data:")
            missing_data = df.isnull().sum().to_frame(name="Missing Values")
            missing_data["% Missing"] = (missing_data["Missing Values"] / len(df)) * 100
            st.dataframe(missing_data.style.background_gradient(cmap='Reds'), 
                        use_container_width=True)
            export_data['tables']['missing_values'] = missing_data
    
    # 2. Key Performance Indicators (KPIs)
    if col_types['numeric']:
        st.header("📈 Key Performance Indicators")
        st.write("Track and visualize your most important metrics over time.")
        
        selected_kpis = st.multiselect(
            "Select KPIs to track (choose up to 4)", 
            col_types['numeric'], 
            default=col_types['numeric'][:min(4, len(col_types['numeric']))],
            help="Select the numerical metrics you want to track as KPIs"
        )
        
        if selected_kpis:
            st.subheader("KPI Summary")
            cols = st.columns(len(selected_kpis))
            for i, kpi in enumerate(selected_kpis):
                with cols[i]:
                    if len(df) > 1 and col_types['date']:
                        try:
                            delta = (df[kpi].iloc[-1] - df[kpi].iloc[0]) / df[kpi].iloc[0] * 100
                        except:
                            delta = 0
                    else:
                        delta = 0
                    
                    # Enhanced metric display
                    with st.container():
                        st.markdown(f"**{kpi.replace('_', ' ').title()}**")
                        st.markdown(f"<h3 style='margin-top:0; margin-bottom:0;'>{df[kpi].mean():,.2f}</h3>", 
                                   unsafe_allow_html=True)
                        if len(df) > 1 and delta != 0:
                            st.markdown(f"<span style='color: {'green' if delta > 0 else 'red'};'>"
                                       f"{'↑' if delta > 0 else '↓'} {abs(delta):.1f}%</span>", 
                                       unsafe_allow_html=True)
            
            # Trend lines for selected KPIs
            if col_types['date'] and len(selected_kpis) > 0:
                st.subheader("KPI Trends Over Time")
                date_col = col_types['date'][0]
                fig = px.line(df, x=date_col, y=selected_kpis, 
                             title="",
                             labels={'value': 'Metric Value', 'variable': 'KPI'},
                             template='plotly_white',
                             height=500)
                
                # Add range slider if there are many data points
                if len(df) > 30:
                    fig.update_layout(
                        xaxis=dict(
                            rangeselector=dict(
                                buttons=list([
                                    dict(count=1, label="1m", step="month", stepmode="backward"),
                                    dict(count=6, label="6m", step="month", stepmode="backward"),
                                    dict(count=1, label="YTD", step="year", stepmode="todate"),
                                    dict(count=1, label="1y", step="year", stepmode="backward"),
                                    dict(step="all")
                                ])
                            ),
                            rangeslider=dict(visible=True),
                            type="date"
                        )
                    )
                
                st.plotly_chart(fig, use_container_width=True)
                export_data['figures'].append(('kpi_trends', fig))
    
    # 3. Profitability and Financial Analysis
    profit_related = [col for col in col_types['numeric'] if 'profit' in col.lower() or 'revenue' in col.lower() or 'margin' in col.lower()]
    if profit_related and len(profit_related) > 0:
        st.header("💰 Profitability Analysis")
        st.write("Analyze your financial metrics and their relationships.")
        
        col1, col2 = st.columns(2)
        with col1:
            selected_profit = st.selectbox(
                "Select financial metric", 
                profit_related,
                help="Choose a financial metric to analyze"
            )
            
            if col_types['date']:
                st.subheader(f"{selected_profit.replace('_', ' ').title()} Over Time")
                fig = px.area(df, x=col_types['date'][0], y=selected_profit, 
                             title="",
                             template='plotly_white',
                             height=400)
                st.plotly_chart(fig, use_container_width=True)
                export_data['figures'].append(('profit_trend', fig))
            else:
                st.subheader(f"Distribution of {selected_profit.replace('_', ' ').title()}")
                fig = px.histogram(df, x=selected_profit, 
                                 title="",
                                 template='plotly_white',
                                 height=400)
                st.plotly_chart(fig, use_container_width=True)
                export_data['figures'].append(('profit_distribution', fig))
        
        with col2:
            if len(profit_related) > 1:
                st.subheader("Financial Metrics Correlation")
                profit_corr = df[profit_related].corr()
                fig = px.imshow(profit_corr, 
                                text_auto=True, 
                                title="",
                                color_continuous_scale='RdYlGn',
                                zmin=-1, zmax=1,
                                height=400)
                st.plotly_chart(fig, use_container_width=True)
                export_data['figures'].append(('profit_correlation', fig))
                export_data['tables']['profit_correlation'] = profit_corr
            else:
                st.info("Add more financial metrics (like revenue, profit, margin) to see correlation analysis.")
    
    # 4. Customer/Product Segmentation with Pie Charts
    if col_types['categorical'] and len(col_types['categorical']) > 0:
        st.header("👥 Segmentation Analysis")
        st.write("Explore how your data breaks down across different categories.")
        
        col1, col2 = st.columns(2)
        with col1:
            seg_col = st.selectbox(
                "Select category to segment by", 
                col_types['categorical'],
                help="Choose a categorical column to analyze"
            )
        
        with col2:
            if col_types['numeric'] and len(col_types['numeric']) > 0:
                value_col = st.selectbox(
                    "Select metric to analyze", 
                    col_types['numeric'],
                    help="Choose a numerical metric to analyze across segments"
                )
        
        if col_types['numeric'] and len(col_types['numeric']) > 0:
            st.subheader("Segmentation Visualization")
            
            # Add visualization type selector with better labels
            chart_type = st.radio(
                "Choose visualization type", 
                ["Bar Chart", "Pie Chart", "Treemap", "Sunburst"],
                horizontal=True,
                help="Select how you want to visualize the segmentation"
            )
            
            seg_data = df.groupby(seg_col)[value_col].sum().sort_values(ascending=False).reset_index()
            export_data['tables']['segmentation_data'] = seg_data
            
            if chart_type == "Bar Chart":
                fig = px.bar(seg_data, 
                            x=seg_col, 
                            y=value_col,
                            title="",
                            labels={seg_col: seg_col.replace('_', ' ').title(), 
                                   value_col: value_col.replace('_', ' ').title()},
                            template='plotly_white',
                            height=500)
                st.plotly_chart(fig, use_container_width=True)
                export_data['figures'].append(('segmentation_bar', fig))
            
            elif chart_type == "Pie Chart":
                col1, col2 = st.columns([3, 1])
                with col1:
                    fig = px.pie(seg_data, 
                                names=seg_col, 
                                values=value_col,
                                title="",
                                hole=0.3,
                                template='plotly_white')
                    st.plotly_chart(fig, use_container_width=True)
                    export_data['figures'].append(('segmentation_pie', fig))
                with col2:
                    st.write("**Top Segments**")
                    st.dataframe(seg_data.head(10), height=400)
            
            elif chart_type == "Treemap":
                fig = px.treemap(seg_data, 
                                path=[seg_col], 
                                values=value_col,
                                title="",
                                template='plotly_white',
                                height=500)
                st.plotly_chart(fig, use_container_width=True)
                export_data['figures'].append(('segmentation_treemap', fig))
            
            elif chart_type == "Sunburst":
                # For sunburst, we need at least two categorical dimensions
                other_cats = [col for col in col_types['categorical'] if col != seg_col]
                if other_cats:
                    second_cat = st.selectbox(
                        "Select second category for hierarchy", 
                        other_cats,
                        help="Choose a second category to create a hierarchical view"
                    )
                    seg_data = df.groupby([seg_col, second_cat])[value_col].sum().reset_index()
                    fig = px.sunburst(seg_data, 
                                    path=[seg_col, second_cat], 
                                    values=value_col,
                                    title="",
                                    template='plotly_white',
                                    height=600)
                    st.plotly_chart(fig, use_container_width=True)
                    export_data['figures'].append(('segmentation_sunburst', fig))
                else:
                    st.warning("Sunburst chart requires at least two categorical columns. Add another category column to your data to use this visualization.")
    
    # 5. Geographic Analysis (if location data exists)
    location_cols = [col for col in df.columns if 'country' in col.lower() or 'state' in col.lower() or 'city' in col.lower()]
    if location_cols and len(location_cols) > 0 and col_types['numeric'] and len(col_types['numeric']) > 0:
        st.header("🌍 Geographic Analysis")
        st.write("Visualize your data geographically when location information is available.")
        
        col1, col2 = st.columns(2)
        with col1:
            loc_col = st.selectbox(
                "Select location column", 
                location_cols,
                help="Choose the column that contains geographic information"
            )
        with col2:
            metric_col = st.selectbox(
                "Select metric to visualize", 
                col_types['numeric'],
                help="Choose a numerical metric to visualize on the map"
            )
        
        loc_data = df.groupby(loc_col)[metric_col].sum().reset_index()
        export_data['tables']['geographic_data'] = loc_data
        
        # Add visualization type selector
        geo_chart_type = st.radio(
            "Select map visualization type", 
            ["Choropleth Map", "Bar Chart", "Pie Chart"],
            horizontal=True,
            help="Choose how to visualize geographic distribution"
        )
        
        if geo_chart_type == "Choropleth Map":
            st.subheader(f"Geographic Distribution of {metric_col}")
            # Try to plot as map if country codes are recognizable
            try:
                fig = px.choropleth(loc_data, 
                                   locations=loc_col,
                                   locationmode="country names",
                                   color=metric_col,
                                   hover_name=loc_col,
                                   title="",
                                   template='plotly_white',
                                   height=600)
                st.plotly_chart(fig, use_container_width=True)
                export_data['figures'].append(('geo_choropleth', fig))
            except:
                st.warning("Could not create map visualization. Showing bar chart instead.")
                fig = px.bar(loc_data.sort_values(metric_col, ascending=False),
                            x=loc_col, y=metric_col,
                            title="",
                            template='plotly_white',
                            height=500)
                st.plotly_chart(fig, use_container_width=True)
                export_data['figures'].append(('geo_bar', fig))
        
        elif geo_chart_type == "Bar Chart":
            st.subheader(f"{metric_col} by {loc_col}")
            fig = px.bar(loc_data.sort_values(metric_col, ascending=False),
                        x=loc_col, y=metric_col,
                        title="",
                        template='plotly_white',
                        height=500)
            st.plotly_chart(fig, use_container_width=True)
            export_data['figures'].append(('geo_bar', fig))
        
        elif geo_chart_type == "Pie Chart":
            st.subheader(f"Distribution of {metric_col} by {loc_col}")
            col1, col2 = st.columns([3, 1])
            with col1:
                fig = px.pie(loc_data, 
                            names=loc_col, 
                            values=metric_col,
                            title="",
                            hole=0.3,
                            template='plotly_white')
                st.plotly_chart(fig, use_container_width=True)
                export_data['figures'].append(('geo_pie', fig))
            with col2:
                st.write("**Top Locations**")
                st.dataframe(loc_data.sort_values(metric_col, ascending=False).head(10), height=400)
    
    # 6. Time Series Forecasting
    if col_types['date'] and len(col_types['date']) > 0 and col_types['numeric'] and len(col_types['numeric']) > 0:
        st.header("🔮 Forecasting")
        st.write("Generate simple forecasts for your time-series data.")
        
        col1, col2 = st.columns(2)
        with col1:
            date_col = st.selectbox(
                "Select date column", 
                col_types['date'],
                help="Choose the column that contains date/time information"
            )
        with col2:
            forecast_col = st.selectbox(
                "Select metric to forecast", 
                col_types['numeric'],
                help="Choose a numerical metric to forecast"
            )
        
        if st.button("Generate Forecast", help="Click to generate a 6-period forecast"):
            with st.spinner("Creating forecast..."):
                try:
                    from statsmodels.tsa.arima.model import ARIMA
                    
                    ts_data = df.set_index(date_col)[forecast_col].dropna()
                    if len(ts_data) > 10:
                        model = ARIMA(ts_data, order=(1,1,1))
                        model_fit = model.fit()
                        forecast = model_fit.forecast(steps=6)  # 6 periods ahead
                        
                        fig = go.Figure()
                        fig.add_trace(go.Scatter(
                            x=ts_data.index,
                            y=ts_data,
                            name="Historical Data",
                            line=dict(width=3)
                        ))
                        fig.add_trace(go.Scatter(
                            x=pd.date_range(ts_data.index[-1], periods=7)[1:],
                            y=forecast,
                            name="Forecast",
                            line=dict(dash='dot', width=3, color='red')
                        ))
                        fig.update_layout(
                            title=f"6-Period Forecast for {forecast_col}",
                            template='plotly_white',
                            height=500,
                            legend=dict(
                                orientation="h",
                                yanchor="bottom",
                                y=1.02,
                                xanchor="right",
                                x=1
                            )
                        )
                        st.plotly_chart(fig, use_container_width=True)
                        export_data['figures'].append(('forecast', fig))
                        
                        # Show forecast values
                        st.subheader("Forecast Values")
                        forecast_df = pd.DataFrame({
                            'Period': range(1, 7),
                            'Date': pd.date_range(ts_data.index[-1], periods=7)[1:],
                            'Forecast': forecast
                        })
                        st.dataframe(forecast_df.set_index('Period'), use_container_width=True)
                        export_data['tables']['forecast_values'] = forecast_df
                    else:
                        st.warning("Not enough data points for reliable forecasting. Need at least 10 observations.")
                except Exception as e:
                    st.error(f"Forecasting failed: {str(e)}")
                    st.info("Try ensuring your time series has enough data points and isn't too irregular.")
    
    # Export functionality
    st.markdown("---")
    st.header("📤 Export Analysis")
    
    col1, col2, col3 = st.columns(3)
    with col1:
        st.markdown("### Excel Export")
        st.write("Download all analysis as an Excel workbook with multiple sheets")
        if st.button("Export to Excel", key="excel_export"):
            with st.spinner("Preparing Excel export..."):
                try:
                    excel_buffer = export_to_excel(export_data)
                    st.success("Excel file ready for download!")
                    st.download_button(
                        label="Download Excel File",
                        data=excel_buffer,
                        file_name="business_analysis.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                    )
                except Exception as e:
                    st.error(f"Error creating Excel file: {str(e)}")
    
    with col2:
        st.markdown("### PDF Export")
        st.write("Create a PDF report with all visualizations and insights")
        if st.button("Export to PDF", key="pdf_export"):
            with st.spinner("Preparing PDF report..."):
                try:
                    pdf_buffer = export_to_pdf(export_data)
                    st.success("PDF report ready for download!")
                    st.download_button(
                        label="Download PDF Report",
                        data=pdf_buffer,
                        file_name="business_analysis.pdf",
                        mime="application/pdf"
                    )
                except Exception as e:
                    st.error(f"Error creating PDF: {str(e)}")
    
    with col3:
        st.markdown("### PowerPoint Export")
        st.write("Generate a PowerPoint presentation with key insights")
        if st.button("Export to PowerPoint", key="ppt_export"):
            with st.spinner("Preparing PowerPoint..."):
                try:
                    ppt_buffer = export_to_ppt(export_data)
                    st.success("PowerPoint ready for download!")
                    st.download_button(
                        label="Download PowerPoint",
                        data=ppt_buffer,
                        file_name="business_analysis.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                except Exception as e:
                    st.error(f"Error creating PowerPoint: {str(e)}")

def export_to_excel(export_data):
    """Export all analysis to an Excel workbook"""
    output = BytesIO()
    writer = pd.ExcelWriter(output, engine='xlsxwriter')
    
    # Write summary data
    summary_df = pd.DataFrame({
        'Metric': ['Total Records', 'Total Features', 'Missing Values', 'Duplicate Rows'],
        'Value': [
            export_data['summary']['num_records'],
            export_data['summary']['num_features'],
            export_data['summary']['missing_values'],
            export_data['summary']['duplicate_rows']
        ]
    })
    summary_df.to_excel(writer, sheet_name='Summary', index=False)
    
    # Write all tables to separate sheets
    for sheet_name, table_data in export_data['tables'].items():
        if isinstance(table_data, pd.DataFrame):
            table_data.to_excel(writer, sheet_name=sheet_name, index=False)
    
    # Save figures as images and add to Excel
    workbook = writer.book
    for i, (fig_name, fig) in enumerate(export_data['figures']):
        # Create a temporary image file
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmpfile:
            fig.write_image(tmpfile.name, scale=2)
            
            # Add worksheet for the figure
            worksheet = workbook.add_worksheet(fig_name[:30])  # Limit sheet name length
            
            # Insert the image
            worksheet.insert_image('A1', tmpfile.name)
            
            # Delete the temporary file
            os.unlink(tmpfile.name)
    
    writer.close()
    return output.getvalue()

def export_to_pdf(export_data):
    """Export all analysis to a PDF report"""
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    
    # Add title page
    pdf.add_page()
    pdf.set_font('Arial', 'B', 16)
    pdf.cell(0, 10, 'Business Data Analysis Report', 0, 1, 'C')
    pdf.ln(10)
    pdf.set_font('Arial', '', 12)
    pdf.cell(0, 10, f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", 0, 1, 'C')
    
    # Add summary section
    pdf.add_page()
    pdf.set_font('Arial', 'B', 14)
    pdf.cell(0, 10, '1. Executive Summary', 0, 1)
    pdf.set_font('Arial', '', 12)
    
    summary_text = f"""
    Total Records: {export_data['summary']['num_records']:,}
    Total Features: {export_data['summary']['num_features']}
    Missing Values: {export_data['summary']['missing_values']}
    Duplicate Rows: {export_data['summary']['duplicate_rows']}
    """
    pdf.multi_cell(0, 10, summary_text)
    
    # Add figures
    pdf.add_page()
    pdf.set_font('Arial', 'B', 14)
    pdf.cell(0, 10, '2. Key Visualizations', 0, 1)
    
    for fig_name, fig in export_data['figures']:
        # Create a temporary image file
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmpfile:
            fig.write_image(tmpfile.name, scale=2)
            
            # Add figure to PDF
            pdf.set_font('Arial', 'B', 12)
            pdf.cell(0, 10, fig_name.replace('_', ' ').title(), 0, 1)
            pdf.image(tmpfile.name, x=10, w=190)
            pdf.ln(5)
            
            # Delete the temporary file
            os.unlink(tmpfile.name)
    
    # Add tables
    pdf.add_page()
    pdf.set_font('Arial', 'B', 14)
    pdf.cell(0, 10, '3. Data Tables', 0, 1)
    
    for table_name, table_data in export_data['tables'].items():
        if isinstance(table_data, pd.DataFrame):
            pdf.set_font('Arial', 'B', 12)
            pdf.cell(0, 10, table_name.replace('_', ' ').title(), 0, 1)
            
            # Convert DataFrame to list of lists for PDF
            pdf.set_font('Arial', '', 10)
            col_widths = [40, 40, 40]  # Adjust as needed
            
            # Headers
            headers = table_data.columns.tolist()
            for i, header in enumerate(headers):
                pdf.cell(col_widths[i], 10, str(header), 1)
            pdf.ln()
            
            # Data rows
            for _, row in table_data.iterrows():
                for i, col in enumerate(headers):
                    pdf.cell(col_widths[i], 10, str(row[col]), 1)
                pdf.ln()
            
            pdf.ln(5)
    
    # Save PDF to buffer
    pdf_buffer = BytesIO()
    pdf_buffer.write(pdf.output(dest='S').encode('latin1'))
    return pdf_buffer.getvalue()

def export_to_ppt(export_data):
    """Export key insights to a PowerPoint presentation"""
    prs = Presentation()
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Business Data Insights"
    subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}"
    
    # Add summary slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = 'Executive Summary'
    
    tf = body_shape.text_frame
    tf.text = 'Key Metrics:'
    
    p = tf.add_paragraph()
    p.text = f"• Total Records: {export_data['summary']['num_records']:,}"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"• Total Features: {export_data['summary']['num_features']}"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"• Missing Values: {export_data['summary']['missing_values']}"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"• Duplicate Rows: {export_data['summary']['duplicate_rows']}"
    p.level = 1
    
    # Add figures (limit to 5 most important ones)
    for fig_name, fig in export_data['figures'][:5]:
        # Create a temporary image file
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmpfile:
            fig.write_image(tmpfile.name, scale=2)
            
            # Add slide for the figure
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add title
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
            tf = txBox.text_frame
            tf.text = fig_name.replace('_', ' ').title()
            
            # Add image
            slide.shapes.add_picture(tmpfile.name, Inches(1), Inches(1.5), Inches(8), Inches(5))
            
            # Delete the temporary file
            os.unlink(tmpfile.name)
    
    # Save PowerPoint to buffer
    ppt_buffer = BytesIO()
    prs.save(ppt_buffer)
    return ppt_buffer.getvalue()

def main():
    st.sidebar.title("Business Insights Dashboard")
    st.sidebar.image("https://cdn-icons-png.flaticon.com/512/2092/2092693.png", width=80)
    st.sidebar.markdown("""
    <style>
        .sidebar .sidebar-content {
            padding-top: 1rem;
        }
    </style>
    """, unsafe_allow_html=True)
    
    st.sidebar.markdown("### Upload Your Data")
    uploaded_file = st.sidebar.file_uploader(
        "Choose a file", 
        type=['csv', 'xlsx', 'json'],
        help="Supported formats: CSV, Excel (xlsx), JSON"
    )
    
    st.sidebar.markdown("---")
    st.sidebar.markdown("### How to Use")
    st.sidebar.info("""
    1. Upload your business data file
    2. Explore the automatic insights
    3. Customize visualizations using the controls
    4. Download any charts or tables as needed
    """)
    
    st.sidebar.markdown("---")
    st.sidebar.markdown("### Supported Data Types")
    st.sidebar.info("""
    - Financial metrics
    - Sales transactions
    - Customer demographics
    - Operational data
    - Marketing performance
    - Time series data
    """)
    
    st.sidebar.markdown("---")
    st.sidebar.markdown("### About")
    st.sidebar.info("""
    This dashboard helps you quickly analyze business data and gain insights without coding.
    """)
    
    if uploaded_file is not None:
        with st.spinner('Analyzing your data...'):
            df = load_data(uploaded_file)
            
        if df is not None:
            # Add business-specific data enhancements
            date_cols = [col for col in df.columns if pd.api.types.is_datetime64_any_dtype(df[col])]
            if date_cols and len(date_cols) > 0:
                try:
                    df['year'] = df[date_cols[0]].dt.year
                    df['quarter'] = df[date_cols[0]].dt.to_period('Q').astype(str)
                    df['month'] = df[date_cols[0]].dt.month_name()
                    df['weekday'] = df[date_cols[0]].dt.day_name()
                except:
                    pass
            
            col_types = detect_column_types(df)
            create_visualizations(df, col_types)
    else:
        st.markdown("""
        <div class="welcome-container">
            <h1>Business Data Insights Dashboard</h1>
            <p style="font-size: 1.2rem; color: #666;">
                Upload your business data to get started with automatic analysis and visualization
            </p>
            <div style="margin: 3rem 0;">
                <img src="https://cdn-icons-png.flaticon.com/512/3713/3713547.png" width="150">
            </div>
            <div class="feature-card">
                <h3 style="margin-top: 0;">Key Features:</h3>
                <ul style="columns: 2; column-gap: 2rem;">
                    <li>📊 Automatic data profiling</li>
                    <li>📈 KPI tracking and visualization</li>
                    <li>💰 Profitability analysis</li>
                    <li>👥 Customer/segment analysis</li>
                    <li>🌍 Geographic visualization</li>
                    <li>🔮 Time series forecasting</li>
                    <li>📋 Interactive data exploration</li>
                    <li>📤 Export to Excel, PDF, PowerPoint</li>
                </ul>
            </div>
        </div>
        """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()
